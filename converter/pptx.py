"""PPTX -> Markdown converter built on python-pptx."""
from __future__ import annotations

import time
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

from converter import config
from converter.base import (
    ConvertResult,
    Converter,
    PageProgressCallback,
    _escape,
    _format_md,
    _link_dest,
    _table_to_md,
    image_digest,
    repeated_image_hashes,
    write_image,
)
from converter.classify import (
    maybe_transcribe_image,
    should_transcribe,
    transcribe_image_cached,
)
from converter.logstore import record
from converter.render import PPTXRenderer, emu_rect_to_points, soffice_available
from converter.vision import VISION_MODEL, transcription_quality

SKIP_PLACEHOLDERS = {PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE}
FOOTER_PLACEHOLDERS = {PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.HEADER}
TITLE_PLACEHOLDERS = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}
BODY_PLACEHOLDERS = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.VERTICAL_BODY,
    PP_PLACEHOLDER.OBJECT,
}
PICTURE_PLACEHOLDERS = {
    PP_PLACEHOLDER.PICTURE,
    PP_PLACEHOLDER.OBJECT,
    PP_PLACEHOLDER.ORG_CHART,
    PP_PLACEHOLDER.BITMAP,
}


def _run_to_md(run) -> str:
    return _format_md(run.text, bool(run.font.bold), bool(run.font.italic))


def _paragraph_style(p) -> tuple[str | None, str | None]:
    """Return (kind, char) where kind is 'bullet' | 'number' | 'none' | None."""
    pPr = p._p.find(qn("a:pPr"))
    if pPr is None:
        return None, None
    if pPr.find(qn("a:buChar")) is not None:
        char = pPr.find(qn("a:buChar")).get("char", "-")
        return "bullet", char
    if pPr.find(qn("a:buAutoNum")) is not None:
        return "number", None
    if pPr.find(qn("a:buNone")) is not None:
        return "none", None
    return None, None


def _paragraph_to_md(p, kind: str | None, default_bullet: bool = False) -> str:
    text = "".join(_run_to_md(run) for run in p.runs)
    if not text.strip():
        return ""
    level = p.level or 0
    if kind == "bullet" or (kind is None and (level > 0 or default_bullet)):
        indent = "  " * level
        return f"{indent}- {text}"
    if kind == "number":
        indent = "  " * level
        return f"{indent}1. {text}"
    return text


def _placeholder_type(shape):
    if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
        return None
    try:
        return shape.placeholder_format.type
    except Exception:
        return None


def _shape_text_to_md(shape, out: list[str], warnings: list[str]):
    if shape.has_table:
        out.extend(_table_to_md(_table_rows(shape.table)))
        return
    if not shape.has_text_frame:
        return
    default_bullet = _placeholder_type(shape) in BODY_PLACEHOLDERS
    text_frame = shape.text_frame
    paragraphs = list(text_frame.paragraphs)
    pending_blank = False
    for p in paragraphs:
        kind, _ = _paragraph_style(p)
        line = _paragraph_to_md(p, kind, default_bullet)
        if line:
            if pending_blank:
                out.append("")
                pending_blank = False
            out.append(line)
        elif not pending_blank:
            pending_blank = True


def _table_rows(table) -> list[list[str]]:
    return [[cell.text for cell in row.cells] for row in table.rows]


def _picture_bytes(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return shape.image
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        if shape.placeholder_format.type in PICTURE_PLACEHOLDERS:
            return shape.image
    return None


def _handle_image(shape, assets_dir: Path, stem: str, counter: list[int], warnings: list[str], dedup: dict[str, str]):
    try:
        image = _picture_bytes(shape)
    except Exception:
        image = None
    if image is None:
        return None
    try:
        blob = image.blob
        ext = image.ext or "bin"
        size = image.size
    except Exception as exc:
        warnings.append(f"Could not read image '{shape.name}': {exc}")
        return None
    filename = write_image(blob, ext, assets_dir, stem, counter, warnings, dedup)
    if filename is None:
        return None
    width, height = size if size else (None, None)
    return filename, image_digest(blob), blob, ext, width, height


def _shape_image_digests(shape) -> set[str]:
    """Collect content digests of every image in a shape (groups recursed)."""
    digests: set[str] = set()
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            digests |= _shape_image_digests(sub)
        return digests
    try:
        image = _picture_bytes(shape)
    except Exception:
        return digests
    if image is None:
        return digests
    try:
        digests.add(image_digest(image.blob))
    except Exception:
        pass
    return digests


def _walk_shape(shape, ctx) -> list[str]:
    """Return markdown lines for a shape; handles groups recursively."""
    out: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            out.extend(_walk_shape(sub, ctx))
        return out
    if shape.has_table:
        return _table_to_md(_table_rows(shape.table))
    image = _handle_image(shape, ctx["assets_dir"], ctx["stem"], ctx["counter"], ctx["warnings"], ctx["dedup"])
    if image:
        filename, digest, blob, ext, width, height = image
        rel = _link_dest(f"assets/{ctx['stem']}/{filename}")
        if digest in ctx["repeated"]:
            if digest in ctx["seen"]:
                return [f"[{shape.name}]({rel})"]
            ctx["seen"].add(digest)
        out = [f"![{shape.name}]({rel})"]
        transcription = maybe_transcribe_image(
            blob,
            ext,
            ctx["warnings"],
            log_ctx={
                "source": ctx["source"],
                "page": ctx["slide_num"],
                "image_ref": shape.name,
                "image_digest": digest,
            },
            width=width,
            height=height,
        )
        if transcription:
            out.extend(transcription.splitlines())
            out.append("")
        return out
    if shape.has_chart:
        return _handle_chart(shape, ctx)
    if shape.has_text_frame:
        _shape_text_to_md(shape, out, ctx["warnings"])
    return out


def _handle_chart(shape, ctx) -> list[str]:
    """Render a chart via LibreOffice, then classify + transcribe it.

    Charts are only handled when both ``VISION_ENABLED`` and
    ``VISION_CLASSIFY_ENABLED`` are on and LibreOffice is available; otherwise
    the chart is skipped with a warning (the pre-existing behaviour).
    """
    if not (config.is_enabled("vision") and config.is_enabled("classify")):
        ctx["warnings"].append(
            f"Chart '{shape.name}' skipped (set VISION_ENABLED and VISION_CLASSIFY_ENABLED to transcribe charts)"
        )
        return []
    renderer = ctx.get("renderer")
    if renderer is None:
        ctx["warnings"].append(f"Chart '{shape.name}' skipped (LibreOffice not available to render it)")
        return []
    slide_index = ctx["slide_num"] - 1
    page_rect = renderer.page_rect(slide_index)
    rect = emu_rect_to_points(
        shape.left,
        shape.top,
        shape.width,
        shape.height,
        ctx["slide_w"],
        ctx["slide_h"],
        page_rect.width,
        page_rect.height,
    )
    png = renderer.render_rect(slide_index, rect)
    if not png:
        ctx["warnings"].append(f"Chart '{shape.name}' skipped (could not render)")
        return []
    filename = write_image(png, "png", ctx["assets_dir"], ctx["stem"], ctx["counter"], ctx["warnings"], ctx["dedup"])
    if filename is None:
        return []
    rel = _link_dest(f"assets/{ctx['stem']}/{filename}")
    out = [f"![{shape.name}]({rel})"]
    log_ctx = {
        "source": ctx["source"],
        "page": ctx["slide_num"],
        "image_ref": f"chart: {shape.name}",
        "image_digest": image_digest(png),
    }
    if should_transcribe(png, "image/png", ctx["warnings"], log_ctx=log_ctx):
        t0 = time.perf_counter()
        try:
            transcription, usage = transcribe_image_cached(png, "image/png", return_usage=True)
            latency_ms = int((time.perf_counter() - t0) * 1000)
        except Exception as exc:
            ctx["warnings"].append(f"Vision transcription failed: {exc}")
            record(
                source=ctx["source"],
                page=ctx["slide_num"],
                image_ref=f"chart: {shape.name}",
                image_digest=image_digest(png),
                stage="transcribe",
                model=VISION_MODEL,
                latency_ms=int((time.perf_counter() - t0) * 1000),
                error=str(exc),
            )
            transcription = None
        if transcription:
            prompt_tokens = generated_tokens = None
            if usage:
                prompt_tokens = usage.get("prompt_tokens")
                generated_tokens = usage.get("completion_tokens", usage.get("generated_tokens"))
            reason = transcription_quality(transcription)
            if reason is not None:
                ctx["warnings"].append(f"Discarding low-value chart transcription ({reason})")
                record(
                    source=ctx["source"],
                    page=ctx["slide_num"],
                    image_ref=f"chart: {shape.name}",
                    image_digest=image_digest(png),
                    stage="transcribe",
                    model=VISION_MODEL,
                    latency_ms=latency_ms,
                    prompt_tokens=prompt_tokens,
                    generated_tokens=generated_tokens,
                    markdown=transcription,
                    error=f"quality gate: {reason}",
                )
            else:
                record(
                    source=ctx["source"],
                    page=ctx["slide_num"],
                    image_ref=f"chart: {shape.name}",
                    image_digest=image_digest(png),
                    stage="transcribe",
                    model=VISION_MODEL,
                    latency_ms=latency_ms,
                    prompt_tokens=prompt_tokens,
                    generated_tokens=generated_tokens,
                    markdown=transcription,
                )
                out.extend(transcription.splitlines())
                out.append("")
    return out


def _shape_has_chart(shape) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return any(_shape_has_chart(sub) for sub in shape.shapes)
    try:
        return bool(shape.has_chart)
    except Exception:
        return False


def _slide_has_chart(slide) -> bool:
    return any(_shape_has_chart(shape) for shape in slide.shapes)


def _is_skipped_placeholder(shape) -> bool:
    return _placeholder_type(shape) in SKIP_PLACEHOLDERS


def _is_footer_placeholder(shape) -> bool:
    return _placeholder_type(shape) in FOOTER_PLACEHOLDERS


def _is_title_placeholder(shape) -> bool:
    return _placeholder_type(shape) in TITLE_PLACEHOLDERS


def _slide_to_md(slide, ctx) -> list[str]:
    lines: list[str] = []
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.text_frame.text.strip():
        title = _escape(title_shape.text_frame.text.strip())
        heading = f"{title} — Slide {ctx['slide_num']}"
    else:
        heading = f"Slide {ctx['slide_num']}"
    lines.append(f"# {heading}")
    lines.append("")
    for shape in slide.shapes:
        if _is_skipped_placeholder(shape):
            continue
        if _is_title_placeholder(shape):
            continue
        body = _walk_shape(shape, ctx)
        if _is_footer_placeholder(shape):
            body = [f"*{_escape(line.strip())}*" for line in body if line.strip()]
        if body:
            lines.extend(body)
            lines.append("")
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            lines.append("> **Notes:**")
            for note_line in notes.splitlines():
                lines.append(f"> {note_line}")
            lines.append("")
    return lines


class PPTXConverter(Converter):
    extensions = (".pptx",)

    def convert(
        self,
        path: Path,
        output_dir: Path,
        progress_callback: PageProgressCallback | None = None,
    ) -> ConvertResult:
        result = ConvertResult(source_path=path)
        try:
            prs = Presentation(path)
            stem = path.stem
            assets_dir = output_dir / "assets" / stem
            assets_dir.mkdir(parents=True, exist_ok=True)
            slide_count = len(prs.slides)
            per_slide: list[set[str]] = []
            for slide in prs.slides:
                digests: set[str] = set()
                for shape in slide.shapes:
                    digests |= _shape_image_digests(shape)
                per_slide.append(digests)
            ctx = {
                "assets_dir": assets_dir,
                "stem": stem,
                "counter": [1],
                "warnings": result.warnings,
                "slide_num": 0,
                "dedup": {},
                "repeated": repeated_image_hashes(per_slide),
                "seen": set(),
                "renderer": None,
                "slide_w": int(prs.slide_width),
                "slide_h": int(prs.slide_height),
                "source": str(path),
            }
            renderer = None
            has_charts = any(_slide_has_chart(slide) for slide in prs.slides)
            if has_charts and config.is_enabled("vision") and config.is_enabled("classify"):
                if soffice_available():
                    try:
                        renderer = PPTXRenderer(path)
                    except Exception as exc:
                        result.warnings.append(f"Could not render charts: {exc}")
                        renderer = None
                else:
                    result.warnings.append(
                        "Charts present but LibreOffice is not installed; "
                        "chart transcription skipped (brew install --cask libreoffice)"
                    )
            ctx["renderer"] = renderer
            try:
                lines: list[str] = []
                for idx, slide in enumerate(prs.slides, start=1):
                    ctx["slide_num"] = idx
                    lines.extend(_slide_to_md(slide, ctx))
                    lines.extend([
                        "",
                        '<div style="page-break-after: always; break-after: page;"></div>',
                        "",
                        "---",
                        "",
                    ])
                    if progress_callback:
                        progress_callback(idx, slide_count, path.name)
            finally:
                if renderer is not None:
                    renderer.close()
            md_path = output_dir / f"{stem}.md"
            md_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")
            result.md_path = md_path
        except Exception as exc:
            result.error = f"{type(exc).__name__}: {exc}"
        return result
