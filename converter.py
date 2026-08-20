"""Core PPTX -> Markdown conversion library.

Pure library, no UI. Kept separate from gui.py so it can be reused from
scripts, a future CLI, or tests.
"""
from __future__ import annotations

import hashlib
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn

ProgressCallback = Callable[[int, int, str], None]

SKIP_PLACEHOLDERS = {PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE}
FOOTER_PLACEHOLDERS = {PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.HEADER}
TITLE_PLACEHOLDERS = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
    PP_PLACEHOLDER.VERTICAL_TITLE,
}
BODY_PLACEHOLDERS = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.VERTICAL_BODY,
    PP_PLACEHOLDER.OBJECT,
}
PICTURE_PLACEHOLDERS = {
    PP_PLACEHOLDER.PICTURE,
    PP_PLACEHOLDER.OBJECT,
    PP_PLACEHOLDER.ORG_CHART,
    PP_PLACEHOLDER.BITMAP,
}

_MD_SPECIALS = str.maketrans({"\\": "\\\\", "`": "\\`", "*": "\\*", "_": "\\_"})


@dataclass
class ConvertResult:
    pptx_path: Path
    md_path: Path | None = None
    error: str | None = None
    warnings: list[str] = field(default_factory=list)


def _escape(text: str) -> str:
    text = text.translate(_MD_SPECIALS)
    if text.startswith("#"):
        text = "\\" + text
    return text


def _run_to_md(run) -> str:
    text = _escape(run.text)
    if not text:
        return ""
    bold = bool(run.font.bold)
    italic = bool(run.font.italic)
    if bold and italic:
        return f"***{text}***"
    if bold:
        return f"**{text}**"
    if italic:
        return f"*{text}*"
    return text


def _paragraph_style(p) -> tuple[str | None, str | None]:
    """Return (kind, char) where kind is 'bullet' | 'number' | 'none' | None."""
    pPr = p._p.find(qn("a:pPr"))
    if pPr is None:
        return None, None
    if pPr.find(qn("a:buChar")) is not None:
        char = pPr.find(qn("a:buChar")).get("char", "-")
        return "bullet", char
    if pPr.find(qn("a:buAutoNum")) is not None:
        return "number", None
    if pPr.find(qn("a:buNone")) is not None:
        return "none", None
    return None, None


def _paragraph_to_md(p, kind: str | None, default_bullet: bool = False) -> str:
    text = "".join(_run_to_md(run) for run in p.runs)
    if not text.strip():
        return ""
    level = p.level or 0
    if kind == "bullet" or (kind is None and (level > 0 or default_bullet)):
        indent = "  " * level
        return f"{indent}- {text}"
    if kind == "number":
        indent = "  " * level
        return f"{indent}1. {text}"
    return text


def _placeholder_type(shape):
    if shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER:
        return None
    try:
        return shape.placeholder_format.type
    except Exception:
        return None


def _shape_text_to_md(shape, out: list[str], warnings: list[str]):
    if shape.has_table:
        out.extend(_table_to_md(shape.table))
        return
    if not shape.has_text_frame:
        return
    default_bullet = _placeholder_type(shape) in BODY_PLACEHOLDERS
    text_frame = shape.text_frame
    paragraphs = list(text_frame.paragraphs)
    pending_blank = False
    for p in paragraphs:
        kind, _ = _paragraph_style(p)
        line = _paragraph_to_md(p, kind, default_bullet)
        if line:
            if pending_blank:
                out.append("")
                pending_blank = False
            out.append(line)
        elif not pending_blank:
            pending_blank = True


def _table_to_md(table) -> list[str]:
    lines: list[str] = []
    rows = list(table.rows)
    if not rows:
        return lines

    def cell_text(cell) -> str:
        return cell.text.replace("|", "\\|").replace("\n", "<br>").strip()

    def render_row(row) -> str:
        return "| " + " | ".join(cell_text(c) for c in row.cells) + " |"

    header = render_row(rows[0])
    separator = "| " + " | ".join("---" for _ in rows[0].cells) + " |"
    lines.extend([header, separator])
    for row in rows[1:]:
        lines.append(render_row(row))
    return lines


def _picture_bytes(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return shape.image
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        if shape.placeholder_format.type in PICTURE_PLACEHOLDERS:
            return shape.image
    return None


def _handle_image(shape, assets_dir: Path, stem: str, counter: list[int], warnings: list[str]):
    try:
        image = _picture_bytes(shape)
    except Exception:
        image = None
    if image is None:
        return None
    try:
        blob = image.blob
        ext = image.ext or "bin"
    except Exception as exc:
        warnings.append(f"Could not read image '{shape.name}': {exc}")
        return None
    digest = hashlib.md5(blob).hexdigest()[:8]
    filename = f"{stem}_{counter[0]:02d}_{digest}.{ext}"
    counter[0] += 1
    (assets_dir / filename).write_bytes(blob)
    return filename


def _walk_shape(shape, ctx) -> list[str]:
    """Return markdown lines for a shape; handles groups recursively."""
    out: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            out.extend(_walk_shape(sub, ctx))
        return out
    if shape.has_table:
        return _table_to_md(shape.table)
    image = _handle_image(shape, ctx["assets_dir"], ctx["stem"], ctx["counter"], ctx["warnings"])
    if image:
        rel = f"assets/{ctx['stem']}/{image}"
        return [f"![{shape.name}]({rel})"]
    if shape.has_chart:
        ctx["warnings"].append(f"Chart '{shape.name}' skipped (charts not supported yet)")
        return []
    if shape.has_text_frame:
        _shape_text_to_md(shape, out, ctx["warnings"])
    return out


def _is_skipped_placeholder(shape) -> bool:
    return _placeholder_type(shape) in SKIP_PLACEHOLDERS


def _is_footer_placeholder(shape) -> bool:
    return _placeholder_type(shape) in FOOTER_PLACEHOLDERS


def _is_title_placeholder(shape) -> bool:
    return _placeholder_type(shape) in TITLE_PLACEHOLDERS


def _slide_to_md(slide, ctx) -> list[str]:
    lines: list[str] = []
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.text_frame.text.strip():
        title = _escape(title_shape.text_frame.text.strip())
    else:
        title = f"Slide {ctx['slide_num']}"
    lines.append(f"# {title}")
    lines.append("")
    for shape in slide.shapes:
        if _is_skipped_placeholder(shape):
            continue
        if _is_title_placeholder(shape):
            continue
        body = _walk_shape(shape, ctx)
        if _is_footer_placeholder(shape):
            body = [f"*{_escape(line.strip())}*" for line in body if line.strip()]
        if body:
            lines.extend(body)
            lines.append("")
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            lines.append("> **Notes:**")
            for note_line in notes.splitlines():
                lines.append(f"> {note_line}")
            lines.append("")
    return lines


def convert_file(
    pptx_path: str | Path,
    output_dir: str | Path,
) -> ConvertResult:
    """Convert one .pptx file to a .md file plus an assets folder."""
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    result = ConvertResult(pptx_path=pptx_path)
    try:
        prs = Presentation(pptx_path)
        stem = pptx_path.stem
        assets_dir = output_dir / "assets" / stem
        assets_dir.mkdir(parents=True, exist_ok=True)
        ctx = {
            "assets_dir": assets_dir,
            "stem": stem,
            "counter": [1],
            "warnings": result.warnings,
            "slide_num": 0,
        }
        lines: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            ctx["slide_num"] = idx
            lines.extend(_slide_to_md(slide, ctx))
        md_path = output_dir / f"{stem}.md"
        md_path.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")
        result.md_path = md_path
    except Exception as exc:
        result.error = f"{type(exc).__name__}: {exc}"
    return result


def convert_files(
    paths: list[str | Path],
    output_dir: str | Path,
    progress_callback: ProgressCallback | None = None,
) -> list[ConvertResult]:
    """Convert multiple .pptx files; errors are captured per file."""
    results: list[ConvertResult] = []
    total = len(paths)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    for idx, path in enumerate(paths, start=1):
        result = convert_file(path, output_dir)
        results.append(result)
        if progress_callback:
            progress_callback(idx, total, Path(path).name)
    return results
