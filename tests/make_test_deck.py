"""Generate a synthetic .pptx exercising every converter feature."""
import struct
import zlib
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


def make_png(path: Path, w=64, h=48, rgb=(120, 160, 220)):
    def chunk(tag, data):
        c = tag + data
        return struct.pack(">I", len(data)) + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    raw = b"".join(b"\x00" + bytes(rgb) * w for _ in range(h))
    path.write_bytes(
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0))
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )


def set_numbered(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"}))


def set_bullet(p, char="•"):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": char}))


def set_no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def main():
    out = Path(__file__).parent / "test_deck.pptx"
    img = Path(__file__).parent / "test_image.png"
    make_png(img)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Deck"
    slide.placeholders[1].text = "Subtitle with *details*"
    slide.notes_slide.notes_text_frame.text = "Cover slide note."

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Bullets and *formatting*"
    body = slide.placeholders[1].text_frame
    p = body.paragraphs[0]
    p.text = "Plain point"
    set_bullet(p)
    p = body.add_paragraph()
    run = p.add_run()
    run.text = "Bold "
    run.font.bold = True
    run = p.add_run()
    run.text = "italic"
    run.font.italic = True
    run = p.add_run()
    run.text = " tail"
    set_bullet(p)
    p = body.add_paragraph()
    p.text = "Nested child"
    p.level = 1
    set_bullet(p, "–")
    p = body.add_paragraph()
    p.text = "Numbered one"
    set_numbered(p)
    p = body.add_paragraph()
    p.text = "Numbered two"
    set_numbered(p)
    p = body.add_paragraph()
    p.text = "Bullets removed explicitly"
    set_no_bullet(p)
    slide.notes_slide.notes_text_frame.text = "First line\nSecond line."

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Table slide"
    table_shape = slide.shapes.add_table(3, 3, Inches(1), Inches(1.5), Inches(10), Inches(3))
    data = [["Name", "Qty", "Note"], ["Alpha", "2", "has | pipe"], ["Beta", "1", "line1\nline2"]]
    for r, row in enumerate(data):
        for c, value in enumerate(row):
            table_shape.table.cell(r, c).text = value

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Image slide"
    slide.shapes.add_picture(str(img), Inches(2), Inches(2), Inches(4), Inches(3))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.text_frame.text = "Slide with no title placeholder."
    slide.notes_slide.notes_text_frame.text = "Untitled slide note."

    prs.save(out)
    print(f"wrote {out}")


if __name__ == "__main__":
    main()
