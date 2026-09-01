"""Smoke tests for both converters and the dispatch API."""
from pathlib import Path

import pytest

from converter import SUPPORTED_EXTENSIONS, convert_file, convert_files

TESTS_DIR = Path(__file__).parent
PPTX = TESTS_DIR / "test_deck.pptx"
PDF = TESTS_DIR / "test_doc.pdf"


def test_supported_extensions():
    assert {".pptx", ".pdf"} <= SUPPORTED_EXTENSIONS


def test_convert_pptx(tmp_path):
    result = convert_file(PPTX, tmp_path)
    assert result.error is None
    assert result.md_path is not None and result.md_path.exists()
    text = result.md_path.read_text(encoding="utf-8")
    assert text.startswith("# Test Deck — Slide 1")
    assert "assets/" in text
    assert "---" in text
    assert 'page-break-after: always' in text


def test_convert_pdf(tmp_path):
    result = convert_file(PDF, tmp_path)
    assert result.error is None
    assert result.md_path is not None and result.md_path.exists()
    text = result.md_path.read_text(encoding="utf-8")
    assert "# Sample PDF Page — Page 1" in text
    assert "# Second Page — Page 2" in text
    assert "Some more body text." in text
    assert "[Page 1](assets/test_doc/" in text
    assert "---" in text
    assert 'page-break-after: always' in text


def test_convert_pdf_bullets(tmp_path):
    result = convert_file(PDF, tmp_path)
    text = result.md_path.read_text(encoding="utf-8")
    assert "- First bullet item" in text
    assert "- Second bullet inline" in text


def test_convert_pdf_bold_formatting(tmp_path):
    result = convert_file(PDF, tmp_path)
    text = result.md_path.read_text(encoding="utf-8")
    assert "**Enterprise Modeling**" in text
    assert "** " not in text


def test_convert_pdf_footer_stripped(tmp_path):
    result = convert_file(PDF, tmp_path)
    text = result.md_path.read_text(encoding="utf-8")
    assert "Sample Footer" not in text


def test_unsupported_extension(tmp_path):
    bogus = tmp_path / "notes.txt"
    bogus.write_text("hello")
    result = convert_file(bogus, tmp_path)
    assert result.error is not None
    assert result.md_path is None


def test_convert_files_mixed(tmp_path):
    results = convert_files([PPTX, PDF], tmp_path)
    assert len(results) == 2
    assert all(r.error is None for r in results)


def test_convert_file_default_output_dir(tmp_path):
    import shutil

    src = tmp_path / "deck.pptx"
    shutil.copy(PPTX, src)
    result = convert_file(src)
    assert result.error is None
    assert result.md_path == tmp_path / "markdown" / "deck.md"
    assert result.md_path.exists()


def test_pptx_link_dest_encodes_spaces(tmp_path):
    import shutil

    src = tmp_path / "My Deck.pptx"
    shutil.copy(PPTX, src)
    result = convert_file(src, tmp_path)
    assert result.error is None
    text = result.md_path.read_text(encoding="utf-8")
    assert "assets/My%20Deck/" in text


def test_pdf_link_dest_encodes_spaces(tmp_path):
    import shutil

    src = tmp_path / "My Doc.pdf"
    shutil.copy(PDF, src)
    result = convert_file(src, tmp_path)
    assert result.error is None
    text = result.md_path.read_text(encoding="utf-8")
    assert "[Page 1](assets/My%20Doc/" in text


def test_pptx_repeated_image(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    from tests.make_test_deck import make_png

    img = tmp_path / "logo.png"
    make_png(img)

    prs = Presentation()
    for i in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Slide {i + 1}"
        if i < 4:
            slide.shapes.add_picture(str(img), Inches(0.5), Inches(0.5), Inches(1), Inches(1))
    src = tmp_path / "deck.pptx"
    prs.save(src)

    result = convert_file(src, tmp_path)
    assert result.error is None
    text = result.md_path.read_text(encoding="utf-8")
    assert text.count("![") == 1
    assert text.count("](") == 4
    assets = list((tmp_path / "assets" / "deck").iterdir())
    assert len(assets) == 1


def test_pdf_repeated_image(tmp_path):
    import pymupdf as fitz

    from tests.make_test_deck import make_png

    img = tmp_path / "logo.png"
    make_png(img)

    doc = fitz.open()
    for i in range(5):
        page = doc.new_page(width=720, height=540)
        page.insert_text((72, 72), f"Page {i + 1}", fontsize=20)
        if i < 4:
            page.insert_image(fitz.Rect(72, 200, 108, 236), filename=str(img))
    src = tmp_path / "doc.pdf"
    doc.save(src)
    doc.close()

    result = convert_file(src, tmp_path)
    assert result.error is None
    text = result.md_path.read_text(encoding="utf-8")
    assert text.count("![image](") == 1
    assert text.count("[image](") - text.count("![image](") == 3


PAPER = TESTS_DIR / "test_paper.pdf"


def test_pdf_paper_mode_structure(tmp_path, monkeypatch):
    monkeypatch.setenv("PDF_MODE", "paper")
    result = convert_file(PAPER, tmp_path)
    assert result.error is None
    text = result.md_path.read_text(encoding="utf-8")

    assert "# What Makes Things Fun to Learn? Heuristics for Design" in text
    assert "*T. Malone · Xerox PARC*" in text
    assert "## Challenge" in text
    assert "## Fantasy" in text
    assert "## Curiosity" in text
    assert "# Page 2" in text
    assert "# Page 3" in text
    assert "A Study of Fun Games" not in text  # running header stripped
    assert "page-break-after" not in text
    assert not any(line.strip() == "---" for line in text.splitlines())


def test_pdf_paper_mode_column_order(tmp_path, monkeypatch):
    monkeypatch.setenv("PDF_MODE", "paper")
    result = convert_file(PAPER, tmp_path)
    assert result.error is None
    text = result.md_path.read_text(encoding="utf-8")

    assert text.index("## Challenge") < text.index("## Fantasy")
    assert (
        text.index("The first challenge is keeping players")
        < text.index("Fantasy can make a game more")
    )


def test_pdf_paper_mode_slide_default(tmp_path):
    result = convert_file(PAPER, tmp_path)
    assert result.error is None
    text = result.md_path.read_text(encoding="utf-8")
    assert "page-break-after" in text


def test_pdf_columns_linearized_in_slide_mode(tmp_path):
    result = convert_file(PAPER, tmp_path)
    assert result.error is None
    text = result.md_path.read_text(encoding="utf-8")
    assert "Raw extracted text" not in text
    assert text.index("The first challenge is keeping players") < text.index(
        "Fantasy can make a game more"
    )


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__]))
